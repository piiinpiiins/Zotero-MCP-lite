"""
Keynote presentation generator.

Creates .pptx files using python-pptx, then optionally converts to .key
using macOS AppleScript + Keynote.app.
"""

import json
import os
import subprocess
import time
from typing import List, Dict, Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ── Slide layout indices (built-in pptx layouts) ──────────────
LAYOUT_TITLE = 0        # Title Slide
LAYOUT_CONTENT = 1      # Title and Content
LAYOUT_SECTION = 2      # Section Header
LAYOUT_BLANK = 6        # Blank


# ── Color scheme ──────────────────────────────────────────────
COLORS = {
    'title_bg': RGBColor(0x1A, 0x1A, 0x2E),      # Dark navy
    'title_text': RGBColor(0xFF, 0xFF, 0xFF),      # White
    'heading': RGBColor(0x1A, 0x1A, 0x2E),         # Dark navy
    'body': RGBColor(0x33, 0x33, 0x33),             # Dark gray
    'accent': RGBColor(0x00, 0x7A, 0xCC),           # Blue accent
    'light_gray': RGBColor(0x66, 0x66, 0x66),       # Light gray for subtitles
}


def _add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    """Add a title slide with large centered title."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])

    # Title
    title_shape = slide.placeholders[0]
    title_shape.text = title
    for para in title_shape.text_frame.paragraphs:
        para.font.size = Pt(40)
        para.font.bold = True
        para.font.color.rgb = COLORS['heading']

    # Subtitle
    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        for para in subtitle_shape.text_frame.paragraphs:
            para.font.size = Pt(20)
            para.font.color.rgb = COLORS['light_gray']

    return slide


def _add_content_slide(prs: Presentation, title: str, body: List[str],
                       notes: str = ""):
    """Add a content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])

    # Title
    title_shape = slide.placeholders[0]
    title_shape.text = title
    for para in title_shape.text_frame.paragraphs:
        para.font.size = Pt(28)
        para.font.bold = True
        para.font.color.rgb = COLORS['heading']

    # Body bullets
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    for i, bullet in enumerate(body):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        # Support indented bullets with "  - " prefix
        if bullet.startswith("  "):
            para.level = 1
            bullet = bullet.strip().lstrip("- ")

        para.text = bullet.lstrip("- ")
        para.font.size = Pt(18)
        para.font.color.rgb = COLORS['body']

    # Speaker notes
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def _add_section_slide(prs: Presentation, title: str, subtitle: str = ""):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION])

    title_shape = slide.placeholders[0]
    title_shape.text = title
    for para in title_shape.text_frame.paragraphs:
        para.font.size = Pt(36)
        para.font.bold = True
        para.font.color.rgb = COLORS['accent']

    if subtitle and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        for para in subtitle_shape.text_frame.paragraphs:
            para.font.size = Pt(18)
            para.font.color.rgb = COLORS['light_gray']

    return slide


def _add_references_slide(prs: Presentation, title: str, references: List[str]):
    """Add a references slide with smaller font."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])

    title_shape = slide.placeholders[0]
    title_shape.text = title
    for para in title_shape.text_frame.paragraphs:
        para.font.size = Pt(28)
        para.font.bold = True
        para.font.color.rgb = COLORS['heading']

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    for i, ref in enumerate(references):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = ref
        para.font.size = Pt(12)
        para.font.color.rgb = COLORS['light_gray']
        para.space_after = Pt(4)

    return slide


# ── Public API ────────────────────────────────────────────────

def create_presentation(slides_data: List[Dict[str, Any]],
                        title: str = "Presentation",
                        output_path: str = "/tmp/presentation.pptx") -> str:
    """
    Create a .pptx presentation from structured slide data.

    Each slide dict should have:
        - type: "title" | "content" | "section" | "references"
        - title: str
        - body: list[str] (bullet points, for content/references)
        - subtitle: str (for title/section slides)
        - notes: str (optional speaker notes)

    Returns the path to the created .pptx file.
    """
    prs = Presentation()

    # Set 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_def in slides_data:
        slide_type = slide_def.get('type', 'content')
        slide_title = slide_def.get('title', '')
        body = slide_def.get('body', [])
        subtitle = slide_def.get('subtitle', '')
        notes = slide_def.get('notes', '')

        if slide_type == 'title':
            _add_title_slide(prs, slide_title, subtitle)
        elif slide_type == 'section':
            _add_section_slide(prs, slide_title, subtitle)
        elif slide_type == 'references':
            _add_references_slide(prs, slide_title, body)
        else:  # content
            _add_content_slide(prs, slide_title, body, notes)

    prs.save(output_path)
    return output_path


def convert_to_keynote(pptx_path: str, key_path: Optional[str] = None) -> str:
    """
    Convert a .pptx file to .key using Keynote via AppleScript.

    Args:
        pptx_path: Path to the .pptx file
        key_path: Output .key path. If None, uses same name with .key extension.

    Returns the path to the .key file.
    """
    pptx_path = os.path.abspath(pptx_path)
    if key_path is None:
        key_path = pptx_path.rsplit('.', 1)[0] + '.key'
    key_path = os.path.abspath(key_path)

    # AppleScript: open .pptx in Keynote, save as .key, then close
    script = f'''
    tell application "Keynote"
        activate
        set theDoc to open POSIX file "{pptx_path}"
        delay 3
        save theDoc in POSIX file "{key_path}"
        close theDoc saving no
    end tell
    '''

    # Write script to temp file to avoid shell quoting issues
    script_file = f"/tmp/_keynote_convert_{int(time.time())}.scpt"
    with open(script_file, 'w') as f:
        f.write(script)

    try:
        subprocess.run(['osascript', script_file],
                       capture_output=True, text=True, timeout=60, check=True)
    except subprocess.TimeoutExpired:
        raise RuntimeError("Keynote conversion timed out (60s). Is Keynote responding?")
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"Keynote conversion failed: {e.stderr}")
    finally:
        if os.path.exists(script_file):
            os.remove(script_file)

    # Verify output exists
    if not os.path.exists(key_path):
        raise RuntimeError(f"Keynote file was not created at {key_path}")

    return key_path


def open_in_keynote(file_path: str):
    """Open a file (.key or .pptx) in Keynote."""
    subprocess.run(['open', '-a', 'Keynote', file_path])
